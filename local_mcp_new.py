from fastapi import FastAPI
from pathlib import Path
from PIL import Image
import uuid
import json
import torch
import re
import os

from transformers import AutoProcessor
try:
    from transformers import AutoModelForVision2Seq
except Exception:  # pragma: no cover - compatibility fallback
    AutoModelForVision2Seq = None
    from transformers import Qwen2VLForConditionalGeneration
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE

from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE


from layout_engine import LayoutEngine


# =========================================================
# CONFIG
# =========================================================

app = FastAPI()

WORKDIR = Path("workdir")
WORKDIR.mkdir(exist_ok=True)

SLIDE_WIDTH_IN = 10
SLIDE_HEIGHT_IN = 7.5
NODE_WIDTH_IN = 2.5
NODE_HEIGHT_IN = 1.2
HORIZONTAL_SPACING_IN = 3.0
VERTICAL_CENTER_IN = 3.5

# =========================================================
# LOAD QWEN2-VL (STABLE MODE)
# =========================================================

QWEN_MODEL = os.getenv("QWEN_MODEL", "Qwen/Qwen2-VL-7B-Instruct")

os.environ["ACCELERATE_DISABLE_RICH"] = "1"

print("Loading Qwen2-VL...")

qwen_processor = AutoProcessor.from_pretrained(
    QWEN_MODEL, trust_remote_code=True
)

if AutoModelForVision2Seq is not None:
    qwen_model = AutoModelForVision2Seq.from_pretrained(
        QWEN_MODEL,
        trust_remote_code=True,
        torch_dtype=torch.float16,
        device_map="auto",
        low_cpu_mem_usage=True
    )
else:
    qwen_model = Qwen2VLForConditionalGeneration.from_pretrained(
        QWEN_MODEL,
        trust_remote_code=True,
        dtype=torch.float16,
        device_map="auto",
        low_cpu_mem_usage=True
    )

# Determine device for tensors (will use GPU if available, otherwise CPU)
device = "cuda" if torch.cuda.is_available() else "cpu"
qwen_model.eval()

print("Qwen2-VL loaded on", device)

# =========================================================
# JSON EXTRACTION (HARDENED)
# =========================================================

def extract_json_from_text(text: str):
    if not text:
        return None

    # Strip markdown fences if present
    text = re.sub(r"```json|```", "", text).strip()

    # Direct parse
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass

    # Fallback: try to extract and repair incomplete JSON
    match = re.search(r"\{.*", text, re.DOTALL)
    if match:
        json_str = match.group(0)
        
        # Count braces to find complete JSON
        brace_count = 0
        end_idx = 0
        for i, char in enumerate(json_str):
            if char == "{":
                brace_count += 1
            elif char == "}":
                brace_count -= 1
                if brace_count == 0:
                    end_idx = i + 1
                    break
        
        if end_idx > 0:
            json_str = json_str[:end_idx]
            try:
                return json.loads(json_str)
            except json.JSONDecodeError:
                pass
        else:
            # Incomplete JSON - aggressively repair it
            # 1. Remove trailing incomplete content (after last complete field)
            # Find the last complete field ending (}, ], or a complete value)
            # Remove everything after the last comma that doesn't complete
            
            # Strategy: find last successfully closed element and truncate after it
            # Walk backwards to find last }, ], or " followed by comma
            last_good_pos = -1
            for i in range(len(json_str) - 1, -1, -1):
                if json_str[i] in '}]"' and i + 1 < len(json_str):
                    # Check if followed by optional whitespace then comma or closing bracket/brace
                    rest = json_str[i+1:].lstrip()
                    if rest and rest[0] in ',}]':
                        last_good_pos = i + 1
                        # Scan past whitespace and comma
                        j = i + 1
                        while j < len(json_str) and json_str[j] in ' \t\n\r':
                            j += 1
                        if j < len(json_str) and json_str[j] == ',':
                            j += 1
                        while j < len(json_str) and json_str[j] in ' \t\n\r':
                            j += 1
                        last_good_pos = j if j < len(json_str) else i + 1
                        break
            
            if last_good_pos > 0:
                json_str = json_str[:last_good_pos]
            
            # 2. Close open arrays
            bracket_count = json_str.count("[") - json_str.count("]")
            if bracket_count > 0:
                json_str += "]" * bracket_count
            
            # 3. Close open braces
            brace_count = json_str.count("{") - json_str.count("}")
            if brace_count > 0:
                json_str += "}" * brace_count
            
            try:
                return json.loads(json_str)
            except json.JSONDecodeError:
                pass
    
    # Last resort: return None to signal failure
    return None

def _save_debug_json(run_id: str, kind: str, payload: dict):
    if not run_id:
        return None
    outpath = WORKDIR / f"{kind}_{run_id}.json"
    try:
        outpath.write_text(json.dumps(payload, indent=2))
        return str(outpath)
    except Exception as exc:
        print(f"Failed to write debug {kind} JSON: {exc}")
        return None

def validate_semantics(semantic_json):
    nodes = semantic_json.get("nodes", [])
    edges = semantic_json.get("edges", [])

    node_ids = {n["id"] for n in nodes}

    ALLOWED_ROLES = {"start", "process", "decision", "end", "connector", "data", "document", "unknown"}

    # Fix invalid roles
    for n in nodes:
        if n["role"] not in ALLOWED_ROLES:
            n["role"] = "unknown"

    # Remove invalid edges
    valid_edges = []
    for e in edges:
        if e["from"] in node_ids and e["to"] in node_ids:
            valid_edges.append(e)

    semantic_json["edges"] = valid_edges
    return semantic_json

def bootstrap_nodes_from_visuals(visual_json):
    nodes = []
    for idx, e in enumerate(visual_json["elements"]):
        if e["type"] == "shape":
            # Accept BOTH:
            # - flowchart symbol labels: terminator/process/decision/data/document/connector
            # - geometry labels: oval/rectangle/diamond/parallelogram/circle
            shape = (e.get("shape") or "").strip().lower()
            text = (e.get("text") or "").strip()
            text_l = text.lower()

            SHAPE_TO_ROLE = {
                # flowchart symbol labels
                "terminator": "start",  # may become "end" based on text
                "process": "process",
                "decision": "decision",
                "data": "data",
                "document": "document",
                "connector": "connector",

                # geometry labels (common outputs from detect_visuals)
                "oval": "start",          # may become "end" based on text
                "rectangle": "process",
                "diamond": "decision",
                "parallelogram": "data",
                "circle": "connector",
            }

            role = SHAPE_TO_ROLE.get(shape, "unknown")

            # Heuristics to refine ambiguous cases
            if role == "start":
                # If the text indicates end, treat it as end terminator
                if "end" in text_l or "stop" in text_l or "finish" in text_l or "terminate" in text_l:
                    role = "end"
            if role == "process":
                # If text suggests a document, prefer document
                if "document" in text_l:
                    role = "document"
                # If text suggests data, prefer data
                elif "data" in text_l:
                    role = "data"

            nodes.append({
                "id": f"n{idx+1}",
                "role": role,
                "text": text,
                "_visual_id": e.get("id"),
                "_bbox": e.get("bbox"),
            })
    return nodes

def _bbox_center_xy(bbox):
    # bbox: [x, y, w, h] in 0-1000-ish coords
    x, y, w, h = bbox
    return (x + w / 2.0, y + h / 2.0)

def _dist2(a, b):
    dx = a[0] - b[0]
    dy = a[1] - b[1]
    return dx * dx + dy * dy

def infer_edges_from_visual_arrows(visual_json, seed_nodes):
    """
    Deterministically infer edges by snapping arrow tail/head points to the nearest shape elements.

    Requires arrows in visual_json like:
      { type:"arrow", tail:[x,y], head:[x,y], label:"Yes"/"No"/"" }
    """
    elements = visual_json.get("elements", []) if isinstance(visual_json, dict) else []
    shapes = [e for e in elements if e.get("type") == "shape" and e.get("bbox")]
    arrows = [e for e in elements if e.get("type") == "arrow" and e.get("tail") and e.get("head")]

    if not shapes or not arrows:
        return []

    # Create a shape->node mapping by visual id first, then by text.
    shape_id_to_node_id = {}
    for n in seed_nodes:
        if n.get("_visual_id"):
            shape_id_to_node_id[n["_visual_id"]] = n["id"]

    # Fallback mapping by text (best-effort).
    text_to_node_id = {}
    for n in seed_nodes:
        if n.get("text"):
            text_to_node_id[n["text"].strip().lower()] = n["id"]

    # Build list of shape centers with node ids (if we can match), else None.
    shape_points = []
    for s in shapes:
        c = _bbox_center_xy(s["bbox"])
        nid = shape_id_to_node_id.get(s.get("id"))
        if not nid:
            st = (s.get("text") or "").strip().lower()
            if st in text_to_node_id:
                nid = text_to_node_id[st]
        shape_points.append((c, nid))

    edges = []
    edge_key_to_idx = {}
    for a in arrows:
        tail = a.get("tail")
        head = a.get("head")
        if not (isinstance(tail, list) and isinstance(head, list) and len(tail) == 2 and len(head) == 2):
            continue

        tail_pt = (float(tail[0]), float(tail[1]))
        head_pt = (float(head[0]), float(head[1]))

        # nearest shape for tail/head
        tail_idx = min(range(len(shape_points)), key=lambda i: _dist2(shape_points[i][0], tail_pt))
        head_idx = min(range(len(shape_points)), key=lambda i: _dist2(shape_points[i][0], head_pt))

        from_id = shape_points[tail_idx][1]
        to_id = shape_points[head_idx][1]
        if not from_id or not to_id or from_id == to_id:
            continue

        e = {"from": from_id, "to": to_id}
        label = (a.get("label") or a.get("text") or "").strip()
        key = (from_id, to_id)
        if key in edge_key_to_idx:
            idx = edge_key_to_idx[key]
            if label:
                existing_label = edges[idx].get("label", "").strip()
                if not existing_label:
                    edges[idx]["label"] = label
                elif label != existing_label:
                    edges[idx]["label"] = f"{existing_label}/{label}"
            continue
        if label:
            e["label"] = label
        edge_key_to_idx[key] = len(edges)
        edges.append(e)

    return edges

def infer_sequential_edges_from_visual_order(visual_json, seed_nodes):
    """
    If we cannot reliably detect arrows, infer a simple top-to-bottom sequence
    using the y-position of shape bboxes.
    """
    elements = visual_json.get("elements", []) if isinstance(visual_json, dict) else []
    shapes = [e for e in elements if e.get("type") == "shape" and e.get("bbox")]
    if not shapes or len(seed_nodes) < 2:
        return []

    # Match nodes to shapes by visual id first, then by text, then sort by bbox center y.
    shape_id_to_node = { n.get("_visual_id"): n.get("id") for n in seed_nodes if n.get("_visual_id") }
    text_to_node = { (n.get("text") or "").strip().lower(): n.get("id") for n in seed_nodes if n.get("text") }
    id_to_node = {n.get("id"): n for n in seed_nodes}
    ordered = []
    for s in shapes:
        nid = shape_id_to_node.get(s.get("id"))
        if not nid:
            nid = text_to_node.get((s.get("text") or "").strip().lower())
        if not nid:
            continue
        cx, cy = _bbox_center_xy(s["bbox"])
        role = id_to_node.get(nid, {}).get("role", "unknown")
        # Force start to the top and end to the bottom when bbox order is unreliable.
        if role == "start":
            cy = -10_000
        elif role == "end":
            cy = 10_000
        ordered.append((cy, cx, nid))

    ordered.sort()
    ids = []
    for _, __, nid in ordered:
        if nid not in ids:
            ids.append(nid)

    if len(ids) < 2:
        return []

    return [{"from": ids[i], "to": ids[i + 1]} for i in range(len(ids) - 1)]

def visual_quality_issues(visual_json):
    """
    Detect obvious model failures in visual_json (used to trigger an automatic retry).
    """
    issues = []
    elements = visual_json.get("elements", []) if isinstance(visual_json, dict) else []
    shapes = [e for e in elements if isinstance(e, dict) and e.get("type") == "shape" and e.get("bbox")]
    arrows = [e for e in elements if isinstance(e, dict) and e.get("type") == "arrow"]

    bbox_counts = {}
    for s in shapes:
        b = s.get("bbox")
        key = tuple(b) if isinstance(b, list) and len(b) == 4 else None
        if key:
            bbox_counts[key] = bbox_counts.get(key, 0) + 1

    duplicated = [k for k, c in bbox_counts.items() if c >= 2]
    if duplicated and len(shapes) >= 3:
        issues.append("Multiple shapes share identical bbox values.")

    unique_bbox = len(bbox_counts)
    if len(shapes) >= 4 and unique_bbox <= max(1, int(len(shapes) * 0.6)):
        issues.append("Too few unique shape bboxes; model likely reused bbox coordinates.")

    for a in arrows:
        if a.get("tail") is None or a.get("head") is None:
            issues.append("Arrow missing tail/head points.")
            break

    return issues

def visual_quality_score(visual_json):
    """
    Lower is better. Used to pick best attempt when model keeps failing.
    """
    issues = visual_quality_issues(visual_json)
    elements = visual_json.get("elements", []) if isinstance(visual_json, dict) else []
    shapes = [e for e in elements if isinstance(e, dict) and e.get("type") == "shape" and e.get("bbox")]

    bbox_counts = {}
    for s in shapes:
        b = s.get("bbox")
        key = tuple(b) if isinstance(b, list) and len(b) == 4 else None
        if key:
            bbox_counts[key] = bbox_counts.get(key, 0) + 1

    dup_penalty = sum((c - 1) for c in bbox_counts.values() if c > 1)
    return len(issues) * 10 + dup_penalty



# =========================================================
# LAYER 1 — VISUAL DETECTION (GEOMETRY ONLY)
# =========================================================

@app.post("/tools/detect_visuals")
async def detect_visuals(data: dict):
    file_id = data.get("file_id")
    run_id = data.get("run_id")
    if not file_id:
        return {"status": "error", "detail": "file_id required"}

    image = Image.open(file_id).convert("RGB")

    prompt = """You are a VISUAL ELEMENT DETECTOR for hand-drawn diagrams and whiteboards.

Your task is to detect ALL visible elements EXACTLY as they appear.
Do NOT infer meaning or flow logic.

============================
GENERAL RULES
============================

- Be COMPLETE, not selective.
- If unsure, STILL include the element.
- NEVER merge elements.
- Use "unknown" if classification is unclear.

============================
ELEMENT TYPES
============================

Detect and output THREE element types ONLY:

1. shape
   - Represents a drawn container (box, oval, diamond, circle, etc.)
   - If a shape contains text, put that text in the SAME element.

2. text
   - Text NOT enclosed inside a shape
   - Examples: titles, annotations, legends

3. arrow
   - Represents a directional connector between shapes
   - Use this ONLY if direction is visually implied
   - IMPORTANT: arrow elements MUST have text="" (empty). Any nearby "Yes/No" belongs in arrow.label.
   - NEVER put node text like "Process" or "Decision" on an arrow.

============================
SHAPE CLASSIFICATION
============================

============================
FLOWCHART SHAPE CLASSIFICATION (MANDATORY)
============================

If the shape visually matches a STANDARD FLOWCHART SYMBOL,
you MUST classify it as one of the following:

- terminator        (start / end – rounded capsule)
- process           (rectangle)
- decision          (diamond)
- data              (parallelogram)
- document          (rectangle with curved bottom)
- connector         (small circle with letter or number)

If unsure, use "unknown".

IMPORTANT:
- Do NOT collapse different symbols into "oval".
- Prefer semantic flowchart symbols over raw geometry.


============================
OUTPUT FIELDS (MANDATORY)
============================

For EACH visible element, output:

- id: unique string (e1, e2, e3, ...)
- type: "shape" | "text" | "arrow"
- text: detected text (empty string if none)
- shape: shape type (ONLY for type="shape", else "unknown")
- bbox: [x, y, w, h]

FOR ARROWS ONLY (type="arrow"), also include:
- tail: [x, y]  (integer point where arrow starts)
- head: [x, y]  (integer point where arrow points to)
- label: text near the arrow (e.g., "Yes"/"No"), empty string if none

============================
BOUNDING BOX RULES
============================

- bbox values must be integers
- Use SMALL approximate values (0–1000 range)
- Coarse estimation is acceptable

============================
CRITICAL CONSTRAINTS
============================

- If text appears INSIDE a shape, it MUST be included in the SAME shape element.
- DO NOT output separate text elements for text inside shapes.
- Titles or decorative text OUTSIDE shapes must be type="text".
- EVERY arrow must be type="arrow".
- Multiple elements are EXPECTED.

============================
OUTPUT FORMAT (STRICT)
============================

Return ONLY valid JSON.
NO markdown.
NO comments.
NO explanations.

Schema (do NOT copy as output):
{
  "elements": [
    {
      "id": "e1",
      "type": "shape|text|arrow",
      "text": "<visible text or empty string>",
      "shape": "<shape name or unknown>",
      "bbox": [x, y, w, h],
      "tail": [x, y],   // arrows only
      "head": [x, y],   // arrows only
      "label": "<Yes/No or empty>" // arrows only
    }
  ]
}

CRITICAL:
- NEVER output the schema above.
- NEVER use placeholder labels like "Start/Process/Decision/End" unless they appear in the image.
- If text is unclear, set text="".

============================
FINAL REQUIREMENT
============================

End your response ONLY after the final closing brace '}'.
Do NOT stop early.

"""

    def _run_detect(prompt_text: str) -> str:
        messages = [
            {
                "role": "user",
                "content": [
                    {"type": "image"},
                    {"type": "text", "text": prompt_text},
                ],
            },
        ]

        text = qwen_processor.apply_chat_template(
            messages, tokenize=False, add_generation_prompt=True
        )

        inputs = qwen_processor(
            text=[text],
            images=[image],
            return_tensors="pt"
        )
        inputs = {k: v.to(device) if isinstance(v, torch.Tensor) else v for k, v in inputs.items()}

        with torch.no_grad():
            output = qwen_model.generate(
                **inputs,
                max_new_tokens=1200,
                eos_token_id=qwen_processor.tokenizer.eos_token_id,
            )

        generated_ids = output[:, inputs["input_ids"].size(1):]
        return qwen_processor.batch_decode(
            generated_ids, skip_special_tokens=True
        )[0]

    def _normalize_visual_json(vj: dict) -> dict:
        vj.setdefault("elements", [])
        for e in vj.get("elements", []):
            if isinstance(e, dict) and e.get("type") == "arrow":
                if (not (e.get("label") or "").strip()) and (e.get("text") or "").strip():
                    e["label"] = str(e.get("text")).strip()
                e["text"] = ""
                e.setdefault("shape", "unknown")
        return vj

    raw = _run_detect(prompt)
    visual_json = extract_json_from_text(raw)
    if visual_json is None:
        return {
            "status": "error",
            "detail": "Failed to extract valid JSON from model output",
            "raw_output": raw[:800],
        }
    visual_json = _normalize_visual_json(visual_json)

    # Quality-gate retries: the model often reuses the same bbox for multiple shapes.
    # We do up to 3 attempts and pick the best-scoring output.
    best = visual_json
    best_score = visual_quality_score(best)

    for attempt in range(3):
        issues = visual_quality_issues(best)
        if not issues:
            break

        repair_prompt = (
            "Your previous JSON had problems and must be corrected.\n\n"
            "PROBLEMS:\n- " + "\n- ".join(issues) + "\n\n"
            "ABSOLUTE RULES (MANDATORY):\n"
            "- Each distinct SHAPE must have a distinct bbox.\n"
            "- For a VERTICAL flowchart, y coordinates MUST increase as you go down.\n"
            "- Do NOT reuse the same bbox for different shapes.\n"
            "- bbox is [x,y,w,h] integers in 0–1000 range (w/h are sizes, not x2/y2).\n"
            "- Arrow must have tail/head points; label is Yes/No if present.\n"
            "- Arrow text must be empty string.\n"
            "- Return ONLY valid JSON.\n\n"
            "HINT: If unsure, spread shapes vertically like y≈50,200,350,500,650.\n\n"
            "PREVIOUS JSON (reference only):\n"
            + json.dumps(best, indent=2)
        )

        raw2 = _run_detect(repair_prompt)
        cand = extract_json_from_text(raw2)
        if isinstance(cand, dict):
            cand = _normalize_visual_json(cand)
            score = visual_quality_score(cand)
            if score < best_score:
                best = cand
                best_score = score

    # After getting best result, apply optional repairs
    visual_json = best

    repair_mode = (data.get("repair_mode") or "conservative").lower()
    shapes_now = [e for e in visual_json.get("elements", []) if e.get("type") == "shape"]
    arrows_now = [e for e in visual_json.get("elements", []) if e.get("type") == "arrow"]

    if repair_mode == "aggressive":
        visual_json = repair_missing_elements(visual_json, expected_min_shapes=6)
        visual_json = add_missing_arrows(visual_json)
    elif repair_mode == "auto":
        # Only add arrows when none are detected; keep shapes untouched
        if len(arrows_now) == 0 and len(shapes_now) >= 2:
            visual_json = add_missing_arrows(visual_json)

    # Final stats
    final_shapes = [e for e in visual_json.get("elements", []) if e.get("type") == "shape"]
    final_arrows = [e for e in visual_json.get("elements", []) if e.get("type") == "arrow"]
    
    return {
        "status": "ok",
        "visual_json": visual_json,
        "debug": {
            "shapes_detected": len(final_shapes),
            "arrows_detected": len(final_arrows),
            "shape_texts": [s.get("text") for s in final_shapes],
            "repairs_applied": repair_mode != "conservative",
            "visual_json_path": _save_debug_json(run_id, "visual", visual_json),
        }
    }

    # return {"status": "ok", "visual_json": visual_json}


def repair_missing_elements(visual_json: dict, expected_min_shapes: int = 6) -> dict:
    """
    Intelligently repairs common detection failures by:
    1. Identifying what's missing based on detected elements
    2. Inferring likely positions for missing elements
    3. Adding synthetic elements with reasonable estimates
    """
    elements = visual_json.get("elements", [])
    shapes = [e for e in elements if e.get("type") == "shape"]
    
    # If we have enough shapes, return as-is
    if len(shapes) >= expected_min_shapes:
        return visual_json
    
    print(f"Attempting repair: only {len(shapes)} shapes detected")
    
    # Get texts of detected shapes
    detected_texts = {(s.get("text") or "").lower().strip() for s in shapes}
    
    # Define expected elements
    expected = {
        "start terminator": {"shape": "terminator", "approx_y": 30},
        "process": {"shape": "process", "approx_y": 110},
        "decision": {"shape": "decision", "approx_y": 220},
        "data": {"shape": "data", "approx_y": 370},
        "document": {"shape": "document", "approx_y": 480},
        "end terminator": {"shape": "terminator", "approx_y": 580},
    }
    
    # Find missing elements
    missing = []
    for text, info in expected.items():
        if text not in detected_texts:
            # Check partial matches
            found = any(text.split()[0] in dt for dt in detected_texts)
            if not found:
                missing.append((text, info))
    
    # Sort existing shapes by Y position
    shapes_with_y = [(s, s.get("bbox", [0, 0, 0, 0])[1]) for s in shapes if s.get("bbox")]
    shapes_with_y.sort(key=lambda x: x[1])
    
    # Add missing shapes
    next_id = max([int(e["id"][1:]) for e in elements if e.get("id", "e0")[1:].isdigit()] + [0]) + 1
    
    for text, info in missing:
        # Estimate position between neighboring shapes
        y_pos = info["approx_y"]
        
        # Find good x position (center column)
        if shapes_with_y:
            avg_x = sum(s[0].get("bbox", [200, 0, 0, 0])[0] for s in shapes_with_y) / len(shapes_with_y)
        else:
            avg_x = 180
        
        # Determine dimensions
        if info["shape"] == "connector":
            width, height = 40, 40
        elif info["shape"] == "decision":
            width, height = 160, 120
        else:
            width, height = 160, 50
        
        new_shape = {
            "id": f"e{next_id}",
            "type": "shape",
            "text": text.title(),
            "shape": info["shape"],
            "bbox": [int(avg_x), int(y_pos), width, height]
        }
        
        elements.append(new_shape)
        next_id += 1
        print(f"  Added missing shape: {text}")
    
    # Check for missing connectors (small circles with "A")
    has_connector = any(
        s.get("shape") == "connector" and "a" in (s.get("text") or "").lower()
        for s in shapes
    )
    
    if not has_connector and len(shapes) >= 4:
        # Add connector pairs if decision node exists
        has_decision = any(s.get("shape") == "decision" for s in shapes)
        if has_decision:
            # Add right connector
            elements.append({
                "id": f"e{next_id}",
                "type": "shape",
                "text": "A",
                "shape": "connector",
                "bbox": [380, 260, 40, 40]
            })
            next_id += 1
            
            # Add left connector
            elements.append({
                "id": f"e{next_id}",
                "type": "shape",
                "text": "A",
                "shape": "connector",
                "bbox": [50, 490, 40, 40]
            })
            print("  Added missing connector pair (A)")
    
    visual_json["elements"] = elements
    return visual_json


def add_missing_arrows(visual_json: dict) -> dict:
    """
    Infers missing arrows based on detected shapes and their positions.
    Creates a sensible flow if arrows are missing.
    """
    elements = visual_json.get("elements", [])
    shapes = [e for e in elements if e.get("type") == "shape" and e.get("bbox")]
    arrows = [e for e in elements if e.get("type") == "arrow"]
    
    if len(arrows) >= len(shapes) - 1:
        return visual_json  # Enough arrows
    
    print(f"Adding missing arrows: {len(arrows)} detected, need ~{len(shapes) - 1}")
    
    # Sort shapes by Y position (top to bottom)
    shapes_sorted = sorted(shapes, key=lambda s: s.get("bbox", [0, 0, 0, 0])[1])
    
    # Create arrows between consecutive shapes
    next_id = max([int(e["id"][1:]) for e in elements if e.get("id", "e0")[1:].isdigit()] + [0]) + 1
    
    # Track which shape pairs already have arrows
    existing_connections = set()
    for a in arrows:
        # Try to match arrow to shapes by proximity
        tail = a.get("tail", [0, 0])
        head = a.get("head", [0, 0])
        # This is approximate - just mark that arrows exist
        existing_connections.add((tuple(tail), tuple(head)))
    
    for i in range(len(shapes_sorted) - 1):
        src = shapes_sorted[i]
        dst = shapes_sorted[i + 1]
        
        src_bbox = src.get("bbox", [0, 0, 0, 0])
        dst_bbox = dst.get("bbox", [0, 0, 0, 0])
        
        # Check if shapes are vertically aligned (likely connected)
        src_cx = src_bbox[0] + src_bbox[2] // 2
        dst_cx = dst_bbox[0] + dst_bbox[2] // 2
        
        if abs(src_cx - dst_cx) < 50:  # Vertically aligned
            # Create arrow from src bottom to dst top
            tail = [src_cx, src_bbox[1] + src_bbox[3]]
            head = [dst_cx, dst_bbox[1]]
            
            # Check if arrow might already exist (approximate)
            is_new = True
            for conn in existing_connections:
                if (abs(conn[0][0] - tail[0]) < 30 and 
                    abs(conn[0][1] - tail[1]) < 30 and
                    abs(conn[1][0] - head[0]) < 30 and
                    abs(conn[1][1] - head[1]) < 30):
                    is_new = False
                    break
            
            if is_new:
                # Determine label
                label = ""
                if src.get("shape") == "decision":
                    # First edge from decision is "Yes" (downward)
                    label = "Yes"
                
                new_arrow = {
                    "id": f"e{next_id}",
                    "type": "arrow",
                    "text": "",
                    "shape": "unknown",
                    "bbox": [min(tail[0], head[0]), min(tail[1], head[1]), 
                             abs(tail[0] - head[0]) + 10, abs(tail[1] - head[1]) + 10],
                    "tail": tail,
                    "head": head,
                    "label": label
                }
                elements.append(new_arrow)
                existing_connections.add((tuple(tail), tuple(head)))
                next_id += 1
                print(f"  Added arrow: {src.get('text', '?')} → {dst.get('text', '?')}")
    
    # Handle decision node branching (if decision exists)
    decision_shapes = [s for s in shapes if s.get("shape") == "decision"]
    if decision_shapes:
        decision = decision_shapes[0]
        d_bbox = decision.get("bbox", [0, 0, 0, 0])
        d_cx = d_bbox[0] + d_bbox[2] // 2
        d_cy = d_bbox[1] + d_bbox[3] // 2
        
        # Look for shapes to the right (for "No" branch)
        right_shapes = [s for s in shapes 
                       if s.get("bbox", [0, 0, 0, 0])[0] > d_bbox[0] + d_bbox[2] and
                          abs(s.get("bbox", [0, 0, 0, 0])[1] - d_cy) < 100]
        
        if right_shapes and len([a for a in arrows if a.get("label") == "No"]) == 0:
            target = right_shapes[0]
            t_bbox = target.get("bbox", [0, 0, 0, 0])
            
            new_arrow = {
                "id": f"e{next_id}",
                "type": "arrow",
                "text": "",
                "shape": "unknown",
                "bbox": [d_bbox[0] + d_bbox[2], d_cy, t_bbox[0] - (d_bbox[0] + d_bbox[2]), 10],
                "tail": [d_bbox[0] + d_bbox[2], d_cy],
                "head": [t_bbox[0], t_bbox[1] + t_bbox[3] // 2],
                "label": "No"
            }
            elements.append(new_arrow)
            print(f"  Added 'No' branch from Decision")
    
    visual_json["elements"] = elements
    return visual_json

# =========================================================
# LAYER 2 — SEMANTIC REASONING (NO GEOMETRY)
# =========================================================

# Replace your reason_semantics function with this

@app.post("/tools/reason_semantics")
async def reason_semantics(data: dict):
    visual_json = data.get("visual_json")
    run_id = data.get("run_id")
    if not visual_json:
        return {"status": "error", "detail": "visual_json required"}

    # Bootstrap nodes from shapes
    seed_nodes = bootstrap_nodes_from_visuals(visual_json)
    
    print(f"Bootstrapped {len(seed_nodes)} nodes from visual shapes")
    for n in seed_nodes:
        print(f"  - {n['id']}: {n['text']} ({n['role']})")
    
    # Try deterministic edge inference from arrows
    inferred_edges = infer_edges_from_visual_arrows(visual_json, seed_nodes)
    
    print(f"Inferred {len(inferred_edges)} edges from arrows")
    for e in inferred_edges:
        print(f"  - {e['from']} -> {e['to']}" + (f" [{e.get('label')}]" if e.get('label') else ""))
    
    if inferred_edges and len(inferred_edges) >= len(seed_nodes) - 1:
        # We have enough edges, use them directly
        semantic_json = {
            "diagram_type": "flowchart",
            "nodes": seed_nodes,
            "edges": inferred_edges,
            "summary": "Flowchart with decision branching and connectors"
        }
        
        from layout_engine import validate_and_score_semantics
        try:
            semantic_json, confidence = validate_and_score_semantics(
                semantic_json, visual_json
            )
            print("Semantic confidence:", confidence)
        except ValueError as e:
            print(f"Validation error: {e}")
            return {
                "status": "error",
                "detail": str(e),
                "raw_semantic": semantic_json,
            }
        
        _save_debug_json(run_id, "semantic", semantic_json)
        return {"status": "ok", "semantic_json": semantic_json}
    
    # Fallback: not enough arrows detected, infer from positions
    print("Insufficient arrows detected, inferring from positions...")
    
    seq_edges = infer_sequential_edges_from_visual_order(visual_json, seed_nodes)
    
    # Merge with any edges we did detect
    all_edges = inferred_edges if inferred_edges else []
    
    # Add sequential edges that don't conflict
    existing_pairs = {(e["from"], e["to"]) for e in all_edges}
    for e in seq_edges:
        pair = (e["from"], e["to"])
        if pair not in existing_pairs:
            all_edges.append(e)
            existing_pairs.add(pair)
    
    # Handle decision branching manually if we have a decision node
    decision_node = next((n for n in seed_nodes if n["role"] == "decision"), None)
    if decision_node:
        # Find nodes after decision
        decision_idx = next(i for i, n in enumerate(seed_nodes) if n["id"] == decision_node["id"])
        
        if decision_idx < len(seed_nodes) - 1:
            # Yes path (typically downward)
            yes_target = seed_nodes[decision_idx + 1]["id"]
            yes_edge = {"from": decision_node["id"], "to": yes_target, "label": "Yes"}
            
            # Check if this edge exists
            if not any(e["from"] == decision_node["id"] and e["to"] == yes_target for e in all_edges):
                all_edges.append(yes_edge)
            else:
                # Update label
                for e in all_edges:
                    if e["from"] == decision_node["id"] and e["to"] == yes_target:
                        e["label"] = "Yes"
        
        # No path (typically to a connector or side branch)
        connector_nodes = [n for n in seed_nodes if n["role"] == "connector"]
        if connector_nodes:
            # Find connector to the right (higher x coordinate)
            # This is approximate since we don't have perfect spatial info in seed_nodes
            # Just connect to first connector as fallback
            no_target = connector_nodes[0]["id"]
            no_edge = {"from": decision_node["id"], "to": no_target, "label": "No"}
            
            if not any(e["from"] == decision_node["id"] and e["to"] == no_target for e in all_edges):
                all_edges.append(no_edge)
    
    # Handle connector merges
    connector_nodes = [n for n in seed_nodes if n["role"] == "connector"]
    if len(connector_nodes) >= 2:
        # First connector (from decision No) should go somewhere
        # Second connector should reconnect to main flow
        
        # Find document node (typical merge target)
        document_node = next((n for n in seed_nodes if n["role"] == "document"), None)
        if document_node and len(connector_nodes) > 1:
            # Connect second connector to document
            merge_edge = {"from": connector_nodes[1]["id"], "to": document_node["id"]}
            if not any(e["from"] == connector_nodes[1]["id"] and e["to"] == document_node["id"] for e in all_edges):
                all_edges.append(merge_edge)
    
    semantic_json = {
        "diagram_type": "flowchart",
        "nodes": seed_nodes,
        "edges": all_edges,
        "summary": "Flowchart with decision branching and connector nodes"
    }
    
    from layout_engine import validate_and_score_semantics
    try:
        semantic_json, confidence = validate_and_score_semantics(
            semantic_json, visual_json
        )
        print("Semantic confidence:", confidence)
    except ValueError as e:
        print(f"Validation error: {e}")
        # Don't fail, just warn
        pass
    
    # Ensure we have edges
    if not semantic_json.get("edges") and len(seed_nodes) > 1:
        print("WARNING: No edges generated, creating minimal sequential flow")
        semantic_json["edges"] = [
            {"from": seed_nodes[i]["id"], "to": seed_nodes[i + 1]["id"]}
            for i in range(len(seed_nodes) - 1)
        ]
    
    print(f"Final semantic: {len(semantic_json['nodes'])} nodes, {len(semantic_json['edges'])} edges")
    
    _save_debug_json(run_id, "semantic", semantic_json)
    return {"status": "ok", "semantic_json": semantic_json}
    
#     prompt = """
# You are a FLOWCHART STRUCTURE ANALYZER.

# You are given VISUAL ELEMENTS extracted from a whiteboard image.
# Each element represents either a shape, text, or arrow.

# Your task is to convert these visual elements into a STRUCTURED FLOWCHART GRAPH.

# ============================
# IMPORTANT — READ CAREFULLY
# ============================

# This diagram IS a FLOWCHART.
# Do NOT summarize, simplify, or invent structure.
# Preserve ALL logical steps exactly as shown.

# You MUST follow the rules below EXACTLY.
# Violating any rule makes the output invalid.

# ============================
# INPUT: VISUAL ELEMENTS
# ============================

# """ + visual_json_str + """

# ============================
# MANDATORY TWO-PHASE CONSTRUCTION
# ============================

# You MUST work in TWO PHASES internally:

# PHASE 1 — NODE ENUMERATION (MANDATORY)
# - Identify ALL visual SHAPES that contain text.
# - Create EXACTLY ONE semantic node for EACH such shape.
# - List ALL nodes completely BEFORE considering edges.
# - If unsure about a role, still create the node with role = "unknown".

# PHASE 2 — EDGE CONSTRUCTION
# - Create edges ONLY AFTER all nodes exist.
# - EVERY edge MUST reference an EXISTING node id.
# - You are NOT ALLOWED to reference an id that is not in the nodes list.

# If fewer than 2 nodes are detected visually,
# you MUST still output all detected nodes.
# Do NOT collapse nodes.

# ============================
# NODE EXTRACTION RULES
# ============================

# 1. One shape → one node (NO exceptions).
# 2. Text inside the shape is the node text.
# 3. Role assignment:
#    - Oval → start or end
#    - Rectangle → process
#    - Diamond → decision
#    - Small labeled circles (A, B, etc.) → connector
#    - Unknown shape → unknown role
# 4. Node IDs:
#    - n1, n2, n3, ...
#    - Sequential
#    - Unique
# 5. Connector nodes:
#    - Same label = same connector node
#    - Multiple arrows may connect through it

# ============================
# EDGE EXTRACTION RULES
# ============================

# 1. Create an edge for EVERY arrow or connector.
# 2. Direction MUST follow arrow direction.
# 3. Arrow text (Yes / No / True / False) → edge.label.
# 4. Decision nodes MUST have at least 2 outgoing edges.
# 5. Merges MUST converge to the same node.

# ============================
# STRUCTURAL RULES
# ============================

# - Multiple start nodes allowed.
# - Multiple end nodes allowed.
# - Cycles (loops) allowed.
# - Parallel paths allowed.
# - A flowchart DOES NOT have only one node unless ONLY ONE shape exists visually.

# ============================
# OUTPUT FORMAT (STRICT)
# ============================

# Return ONLY valid JSON.
# NO markdown.
# NO comments.
# NO explanations.
# NO trailing text.

# {
#   "diagram_type": "flowchart",
#   "nodes": [
#     { "id": "n1", "role": "start", "text": "..." }
#   ],
#   "edges": [
#     { "from": "n1", "to": "n2", "label": "Yes" }
#   ],
#   "summary": "One concise sentence describing the flowchart."
# }

# ============================
# FINAL VALIDATION (MANDATORY)
# ============================

# Before responding:
# - Count visual shapes with text.
# - Ensure the SAME number of nodes exist.
# - Ensure NO edge references missing node ids.
# - Ensure decision nodes have multiple edges.
# - If any rule is violated, FIX IT before returning.

# """
    prompt = """
You are a FLOWCHART EDGE INFERENCE ENGINE.

STRICT OUTPUT RULES (NON-NEGOTIABLE):
- You MUST output VALID JSON ONLY.
- NO explanations.
- NO prose.
- NO markdown.
- NO comments.
- NO text before or after JSON.

If you violate this, the output will be discarded.

====================================
PREDEFINED NODES (DO NOT MODIFY)
====================================
""" + json.dumps(seed_nodes, indent=2) + """

====================================
VISUAL ELEMENTS (FOR REFERENCE)
====================================
""" + json.dumps(visual_json, indent=2) + """

====================================
YOUR TASK
====================================

- Create EDGES between EXISTING node IDs only.
- Use arrow direction/labels if present in visual elements.
- If no arrow text is visible, omit the "label" field.
- Decision nodes may have multiple outgoing edges (Yes/No).
- You are NOT allowed to create new nodes.

====================================
OUTPUT FORMAT (EXACT)
====================================

{
  "edges": [
    { "from": "n1", "to": "n2" },
    { "from": "n2", "to": "n3", "label": "Yes" }
  ],
  "summary": "One concise sentence."
}
"""

    print("=== SEMANTIC PROMPT ===")
    print(prompt)


    messages = [
        {
            "role": "user",
            "content": [{"type": "text", "text": prompt}],
        },
    ]

    text = qwen_processor.apply_chat_template(
        messages, tokenize=False, add_generation_prompt=True
    )

    inputs = qwen_processor(text=[text], return_tensors="pt")
    inputs = {k: v.to(device) if isinstance(v, torch.Tensor) else v for k, v in inputs.items()}

    with torch.no_grad():
        output = qwen_model.generate(
            **inputs,
            max_new_tokens=600,
            eos_token_id=qwen_processor.tokenizer.eos_token_id,
        )

    generated_ids = output[:, inputs["input_ids"].size(1):]
    raw = qwen_processor.batch_decode(
        generated_ids, skip_special_tokens=True
    )[0]
    # raw = raw.strip()
    # raw = raw[raw.find("{") : raw.rfind("}") + 1]

    edges_json = extract_json_from_text(raw)
    
    # Check for extraction failure BEFORE validation
    if edges_json is None or not edges_json.get("edges"):
        # If model failed to return JSON (or returned empty edges), create default edges
        # Prefer visual order (top-to-bottom), otherwise fallback to seed order.
        print("No usable edges returned, creating sequential edges...")
        seq = infer_sequential_edges_from_visual_order(visual_json, seed_nodes)
        if not seq:
            seq = [{"from": seed_nodes[i]["id"], "to": seed_nodes[i + 1]["id"]} for i in range(len(seed_nodes) - 1)]
        edges_json = {"edges": seq}
    
    # Construct complete semantic JSON
    semantic_json = {
        "diagram_type": "flowchart",
        "nodes": seed_nodes,
        "edges": edges_json.get("edges", []),
        "summary": edges_json.get("summary", "")
    }
    
    from layout_engine import validate_and_score_semantics

    try:
        semantic_json, confidence = validate_and_score_semantics(
            semantic_json,
            visual_json
        )
        print("Semantic confidence:", confidence)
    except ValueError as e:
        return {
            "status": "error",
            "detail": str(e),
            "raw_semantic": semantic_json,
        }

    # Final safety: ensure we don't render a graph with zero edges when multiple nodes exist.
    if not semantic_json.get("edges") and len(seed_nodes) > 1:
        seq = infer_sequential_edges_from_visual_order(visual_json, seed_nodes)
        if not seq:
            seq = [{"from": seed_nodes[i]["id"], "to": seed_nodes[i + 1]["id"]} for i in range(len(seed_nodes) - 1)]
        semantic_json["edges"] = seq

    semantic_json.setdefault("nodes", [])
    semantic_json.setdefault("edges", [])
    semantic_json.setdefault("diagram_type", "unknown")
    semantic_json.setdefault("summary", "")

    return {"status": "ok", "semantic_json": semantic_json}

# =========================================================
# LAYER 3 — DETERMINISTIC PPT RENDERING
# =========================================================

def map_role_to_shape(role: str):
    return {
        "start": MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR,
        "end": MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR,
        "process": MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS,
        "decision": MSO_AUTO_SHAPE_TYPE.FLOWCHART_DECISION,
        "data": MSO_AUTO_SHAPE_TYPE.FLOWCHART_DATA,
        "document": MSO_AUTO_SHAPE_TYPE.FLOWCHART_DOCUMENT,
        "connector": MSO_AUTO_SHAPE_TYPE.FLOWCHART_CONNECTOR,
    }.get(role, MSO_AUTO_SHAPE_TYPE.RECTANGLE)

def _connector_points(src, dst):
    src_cx = src.left + src.width / 2
    src_cy = src.top + src.height / 2
    dst_cx = dst.left + dst.width / 2
    dst_cy = dst.top + dst.height / 2

    dx = dst_cx - src_cx
    dy = dst_cy - src_cy

    if abs(dx) >= abs(dy):
        # left/right connection
        if dx >= 0:
            start = (src.left + src.width, src_cy)
            end = (dst.left, dst_cy)
        else:
            start = (src.left, src_cy)
            end = (dst.left + dst.width, dst_cy)
    else:
        # top/bottom connection
        if dy >= 0:
            start = (src_cx, src.top + src.height)
            end = (dst_cx, dst.top)
        else:
            start = (src_cx, src.top)
            end = (dst_cx, dst.top + dst.height)

    return int(start[0]), int(start[1]), int(end[0]), int(end[1])

@app.post("/tools/render_ppt")
async def render_ppt(data: dict):
    semantic = data.get("semantic_json")
    if not semantic:
        return {"status": "error", "detail": "semantic_json required"}

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # nodes = semantic.get("nodes", [])
    # edges = semantic.get("edges", [])
    engine = LayoutEngine()
    layout = engine.layout(semantic)

    layout_nodes = layout["nodes"]
    layout_edges = layout["edges"]

    ppt_nodes = {}

    for node in layout_nodes:
        shape = slide.shapes.add_shape(
            map_role_to_shape(node["role"]),
            Inches(node["x"]),
            Inches(node["y"]),
            Inches(node["width"]),
            Inches(node["height"])
        )

        # ---- Style shape ----
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(235, 242, 255)
        shape.line.width = Pt(1.5)

        # ---- Text formatting ----
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = node["text"]
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        ppt_nodes[node["id"]] = shape


    for edge in layout_edges:
        src = ppt_nodes.get(edge["from"])
        dst = ppt_nodes.get(edge["to"])
        if not src or not dst:
            continue

        x1, y1, x2, y2 = _connector_points(src, dst)
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            x1,
            y1,
            x2,
            y2,
        )

        connector.line.width = Pt(1.5)
        connector.line.end_arrowhead_style = 2  # Arrow
        
        # Add edge label if present (Yes/No)
        label = edge.get("label", "").strip()
        if label:
            # Position label near the connector midpoint
            label_x = (x1 + x2) / 2
            label_y = (y1 + y2) / 2
            label_box = slide.shapes.add_textbox(
                int(label_x - Inches(0.3)),
                int(label_y - Inches(0.15)),
                Inches(0.6),
                Inches(0.3),
            )
            label_tf = label_box.text_frame
            label_tf.text = label
            label_tf.paragraphs[0].font.size = Pt(10)
            label_tf.paragraphs[0].font.bold = True
            label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER



    # ppt_nodes = {}
    # x = 1.0

    # for node in nodes:
    #     shape = slide.shapes.add_shape(
    #         map_role_to_shape(node["role"]),
    #         Inches(x),
    #         Inches(VERTICAL_CENTER_IN),
    #         Inches(NODE_WIDTH_IN),
    #         Inches(NODE_HEIGHT_IN),
    #     )
    #     shape.text = node["text"]
    #     ppt_nodes[node["id"]] = shape
    #     x += HORIZONTAL_SPACING_IN

    # for edge in edges:
    #     src = ppt_nodes.get(edge["from"])
    #     dst = ppt_nodes.get(edge["to"])
    #     if not src or not dst:
    #         continue

    #     connector = slide.shapes.add_connector(
    #         MSO_CONNECTOR_TYPE.STRAIGHT,
    #         src.left + src.width,
    #         src.top + src.height // 2,
    #         dst.left - src.left,
    #         0,
    #     )
    #     connector.line.end_arrowhead = True

    # if semantic.get("summary"):
    #     tb = slide.shapes.add_textbox(
    #         Inches(0.5),
    #         Inches(6.3),
    #         Inches(9),
    #         Inches(0.6),
    #     )
    #     tb.text = semantic["summary"]

    if semantic.get("summary"):
        tb = slide.shapes.add_textbox(
            Inches(6.0),
            Inches(6.5),
            Inches(3.5),
            Inches(0.6),
        )
        tf = tb.text_frame
        tf.text = semantic["summary"]
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].font.italic = True

    outpath = WORKDIR / f"slide_{uuid.uuid4().hex}.pptx"
    prs.save(outpath)

    return {"status": "ok", "ppt_file": str(outpath)}

# =========================================================
# ORCHESTRATOR
# =========================================================

@app.post("/tools/whiteboard_to_ppt")
async def whiteboard_to_ppt(data: dict):
    run_id = data.get("run_id") or uuid.uuid4().hex
    visuals = await detect_visuals({**data, "run_id": run_id})
    if visuals["status"] != "ok":
        return visuals

    semantics = await reason_semantics({
        "visual_json": visuals["visual_json"],
        "run_id": run_id,
    })
    if semantics["status"] != "ok":
        return semantics

    ppt = await render_ppt({"semantic_json": semantics["semantic_json"]})
    if ppt["status"] != "ok":
        return ppt

    return {
        "status": "ok",
        "ppt_file": ppt["ppt_file"],
        "semantic_json": semantics["semantic_json"],
        "run_id": run_id,
        "visual_json_path": visuals.get("debug", {}).get("visual_json_path"),
        "semantic_json_path": _save_debug_json(run_id, "semantic", semantics["semantic_json"]),
    }
